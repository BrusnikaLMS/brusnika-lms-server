import json
import uuid
import asyncio

from fastapi import APIRouter, HTTPException, UploadFile, File, Form, Depends

from app.ai_client import get_ai_client, get_ai_model
from app.config import settings
from app.auth import get_current_user
from app.models import User

router = APIRouter(prefix="/documents", tags=["documents"])


def _extract_docx(data: bytes) -> str:
    import io
    from docx import Document
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(data: bytes) -> str:
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_pdf(data: bytes) -> str:
    import io
    import pdfplumber
    pages = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                pages.append(text.strip())
    return "\n\n".join(pages)


@router.post("/import")
async def import_document(
    file: UploadFile = File(...),
    language: str = Form("en"),
    current_user: User = Depends(get_current_user),
):
    if current_user.plan not in {"pro", "enterprise", "admin"}:
        raise HTTPException(status_code=403, detail="Document import requires Pro plan")

    filename = (file.filename or "").lower()
    if filename.endswith(".docx"):
        extractor = _extract_docx
    elif filename.endswith(".pptx"):
        extractor = _extract_pptx
    elif filename.endswith(".pdf"):
        extractor = _extract_pdf
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type. Use .docx, .pptx, or .pdf")

    data = await file.read()
    text = await asyncio.get_event_loop().run_in_executor(None, extractor, data)

    if not text or not text.strip():
        raise HTTPException(status_code=400, detail="Could not extract text from document")

    text = text[:12000]

    timestamp = str(int(__import__("time").time()))
    system_prompt = "You are a course author. Convert the provided document content into a structured interactive course in JSON format."
    user_prompt = (
        f"Convert this document into an interactive course.\n\n"
        f"Document content:\n{text}\n\n"
        f"Requirements:\n"
        f"- Create a course with a meaningful title based on the document\n"
        f"- Generate 5-15 screens based on the document sections/topics\n"
        f"- Each screen has a title and components array\n"
        f"- Components can be:\n"
        f'  - text: {{"type":"text","id":"c1","content":"<p>...</p>"}}\n'
        f'  - quiz-single: {{"type":"quiz-single","id":"c2","question":"...","options":["a","b","c","d"],"correct":0}}\n'
        f'  - true-false: {{"type":"true-false","id":"c3","statement":"...","correct":true}}\n'
        f"- Include at least one quiz component per 3 screens\n"
        f"- Generate all text in language: {language}\n"
        f"- Return ONLY valid JSON, no markdown, no explanation\n"
        f"- Use this exact structure:\n"
        f'{{"id":"doc-{timestamp}","version":"1.0.0","title":"...","language":"{language}",'
        f'"settings":{{"showResults":true,"allowNavigation":"linear","completionCriteria":"allScreens"}},'
        f'"state":{{}},'
        f'"screens":['
        f'{{"id":"s1","title":"...","components":[...],"events":[],"navigation":{{}}}}'
        f"]}}\n"
        f"- Each component must have a unique id like \"c1\", \"c2\", etc."
    )

    client = get_ai_client()

    try:
        response = await client.chat.completions.create(
            model=get_ai_model(),
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            response_format={"type": "json_object"},
            temperature=0.7,
        )
        raw = response.choices[0].message.content or "{}"

        raw = raw.strip()
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[-1].rsplit("```", 1)[0].strip()

        course = json.loads(raw)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse AI response as JSON: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    if not course.get("id"):
        course["id"] = f"doc-{str(uuid.uuid4())[:8]}"
    for screen in course.get("screens", []):
        if not screen.get("id"):
            screen["id"] = str(uuid.uuid4())[:8]
        for component in screen.get("components", []):
            if not component.get("id"):
                component["id"] = str(uuid.uuid4())[:8]

    return course


@router.post("/transcribe")
async def transcribe_audio(
    file: UploadFile = File(...),
    language: str = Form("en"),
    current_user: User = Depends(get_current_user),
):
    if current_user.plan not in {"pro", "enterprise", "admin"}:
        raise HTTPException(status_code=403, detail="Audio transcription requires Pro plan")

    filename = (file.filename or "").lower()
    allowed_exts = (".mp3", ".mp4", ".wav", ".m4a", ".webm")
    if not any(filename.endswith(ext) for ext in allowed_exts):
        raise HTTPException(status_code=400, detail="Unsupported file type. Use .mp3, .mp4, .wav, .m4a, or .webm")

    data = await file.read()
    if len(data) > 25 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="File too large. Maximum size is 25 MB")

    if not settings.openai_api_key:
        raise HTTPException(status_code=500, detail="OpenAI API key not configured. Set OPENAI_API_KEY in environment.")

    # Transcribe via OpenAI Whisper
    whisper_client = AsyncOpenAI(api_key=settings.openai_api_key)
    try:
        import io
        file_tuple = (filename, io.BytesIO(data), file.content_type or "application/octet-stream")
        transcription = await whisper_client.audio.transcriptions.create(
            model="whisper-1",
            file=file_tuple,
        )
        text = transcription.text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Transcription failed: {e}")

    if not text or not text.strip():
        raise HTTPException(status_code=400, detail="Could not transcribe audio — no speech detected")

    text = text[:12000]

    timestamp = str(int(__import__("time").time()))
    system_prompt = "You are a course author. Convert the provided transcription into a structured interactive course in JSON format."
    user_prompt = (
        f"Convert this audio transcription into an interactive course.\n\n"
        f"Transcription:\n{text}\n\n"
        f"Requirements:\n"
        f"- Create a course with a meaningful title based on the content\n"
        f"- Generate 5-15 screens based on the topics discussed\n"
        f"- Each screen has a title and components array\n"
        f"- Components can be:\n"
        f'  - text: {{"type":"text","id":"c1","content":"<p>...</p>"}}\n'
        f'  - quiz-single: {{"type":"quiz-single","id":"c2","question":"...","options":["a","b","c","d"],"correct":0}}\n'
        f'  - true-false: {{"type":"true-false","id":"c3","statement":"...","correct":true}}\n'
        f"- Include at least one quiz component per 3 screens\n"
        f"- Generate all text in language: {language}\n"
        f"- Return ONLY valid JSON, no markdown, no explanation\n"
        f"- Use this exact structure:\n"
        f'{{"id":"audio-{timestamp}","version":"1.0.0","title":"...","language":"{language}",'
        f'"settings":{{"showResults":true,"allowNavigation":"linear","completionCriteria":"allScreens"}},'
        f'"state":{{}},'
        f'"screens":['
        f'{{"id":"s1","title":"...","components":[...],"events":[],"navigation":{{}}}}'
        f"]}}\n"
        f"- Each component must have a unique id like \"c1\", \"c2\", etc."
    )

    deepseek_client = AsyncOpenAI(
        api_key=settings.deepseek_api_key,
        base_url=settings.deepseek_base_url,
    )

    try:
        response = await deepseek_client.chat.completions.create(
            model=get_ai_model(),
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            response_format={"type": "json_object"},
            temperature=0.7,
        )
        raw = response.choices[0].message.content or "{}"

        raw = raw.strip()
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[-1].rsplit("```", 1)[0].strip()

        course = json.loads(raw)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse AI response as JSON: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    if not course.get("id"):
        course["id"] = f"audio-{str(uuid.uuid4())[:8]}"
    for screen in course.get("screens", []):
        if not screen.get("id"):
            screen["id"] = str(uuid.uuid4())[:8]
        for component in screen.get("components", []):
            if not component.get("id"):
                component["id"] = str(uuid.uuid4())[:8]

    return course
